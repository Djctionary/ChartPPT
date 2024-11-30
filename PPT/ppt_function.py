import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR
import os
import random
from difflib import SequenceMatcher

from Code.FLask.SHAPE.shape_functions import add_basic_shape


def similarity(a, b):
    """

    :param a:
    :param b:
    :return:
    """
    return SequenceMatcher(None, a, b).ratio()


def get_color_string(merged_item):
    """

    :param merged_item:
    :return:
    """
    # 尝试从 merged_item 中获取颜色
    color = merged_item.get('colour')
    if color:
        r, g, b = map(int, color.strip("RGB()").split(", "))
        if 0 <= r <= 255 and 0 <= g <= 255 and 0 <= b <= 255:
            # print("COLOR:", color)
            return color  # 如果找到颜色，直接返回

    # 随机生成 RGB 色值
    r = random.randint(0, 255)
    g = random.randint(0, 255)
    b = random.randint(0, 255)

    # 返回随机生成的颜色字符串
    return f"RGB({r}, {g}, {b})"


def get_shape_string(merged_item):
    """

    :param merged_item:
    :return:
    """
    # 尝试从 merged_item 中获取形状
    shape = merged_item.get('shape')

    if shape:
        return shape  # 如果找到形状，直接返回

    # 返回默认形状字符串
    return "ROUNDED_RECTANGLE"


def add_rounded_rectangle(slide, left, top, width, height, text, color, font_size=14):
    """
    添加一个圆角矩形，并在其中设置文本。

    :param slide: 幻灯片对象
    :param left: 矩形左边距
    :param top: 矩形上边距
    :param width: 矩形宽度
    :param height: 矩形高度
    :param text: 要显示的文本，可以包含换行符
    :param color: 矩形的填充颜色，格式为 "RGB(r, g, b)"
    :param font_size: 字体大小
    :return: 添加的形状对象
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    r, g, b = map(int, color.strip("RGB()").split(", "))
    shape.fill.fore_color.rgb = RGBColor(r, g, b)

    # 设置文本框
    text_frame = shape.text_frame

    # 清空现有文本框内容
    text_frame.clear()  # 确保文本框清空，避免累积文本

    # 分割文本以处理换行
    lines = text.split('\n')

    # 处理第一行文本，直接在默认段落中设置
    if lines:
        p = text_frame.paragraphs[0]  # 获取第一个段落
        p.text = lines[0]  # 设置第一行文本
        p.font.size = Pt(font_size)  # 设置字体大小

        # 设置字体颜色
        if r == g == b == 0:  # 如果颜色是黑色
            font_color = "RGB(255, 255, 255)"  # 设置为白色字体
        else:
            font_color = "RGB(0, 0, 0)"  # 设置为黑色字体
        r_font, g_font, b_font = map(int, font_color.strip("RGB()").split(", "))
        p.font.color.rgb = RGBColor(r_font, g_font, b_font)

    # 处理后续行
    for line in lines[1:]:  # 从第二行开始
        p = text_frame.add_paragraph()  # 添加新的段落
        p.text = line  # 设置每一行文本
        p.font.size = Pt(font_size)  # 设置字体大小

        # 设置字体颜色
        if r == g == b == 0:  # 如果颜色是黑色
            font_color = "RGB(255, 255, 255)"  # 设置为白色字体
        else:
            font_color = "RGB(0, 0, 0)"  # 设置为黑色字体
        r_font, g_font, b_font = map(int, font_color.strip("RGB()").split(", "))
        p.font.color.rgb = RGBColor(r_font, g_font, b_font)

    # 设置文本框垂直居中
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER  # 设置段落居中

    return shape


def get_center(shape):
    """

    :param shape:
    :return:
    """
    center_x = shape.left + shape.width / 2
    center_y = shape.top + shape.height / 2
    return center_x, center_y


def ppt_process(ocr_results, physical_results, logical_results, parameters):
    from Code.FLask.app import app
    ocr_params = parameters.get("ocr_parameters", {})
    ppt_params = parameters.get("ppt_parameters", {})
    user_tag = parameters.get("user_tag")

    recognize_shape = ocr_params.get('recognize_shape')
    # 获取 PPT 参数
    shape_size_ratio = ppt_params.get('shape_size_ratio')
    add_lines_between_shapes = ppt_params.get('add_lines_between_shapes')
    line_style = ppt_params.get('line_style')

    # print(f"Shape Size Ratio: {shape_size_ratio}")
    # print(f"Add Lines Between Shapes: {add_lines_between_shapes}")
    # print(f"Line Style: {line_style}")

    PPT_width = 13.33
    PPT_height = 7.5
    lineSpacing = 0
    # Create a new PowerPoint presentation
    prs = Presentation()
    prs.slide_width = Inches(PPT_width)
    prs.slide_height = Inches(PPT_height)
    slide_layout = prs.slide_layouts[6]  # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Add rounded rectangles for each OCR result
    shapes = []

    merged_results = []
    matched_gpt_indices = set()  # 跟踪已匹配的 GPT 结果

    if not recognize_shape and not add_lines_between_shapes:
        app.logger.info("SKIP Process GPT")
        for ocr_item in ocr_results:
            Inches_box = ocr_item['Inches_box']
            text = ocr_item['text']
            color = get_color_string(ocr_item)
            shape_name = "Rounded Rectangle"

            left = Inches(Inches_box[0][0])
            top = Inches(Inches_box[0][1])
            width = Inches((Inches_box[2][0] - Inches_box[0][0]))
            height = Inches((Inches_box[2][1] - Inches_box[0][1]))

            # print(f"Left: {left.inches:.2f} inches")
            # print(f"Top: {top.inches:.2f} inches")
            # print(f"Width: {width.inches:.2f} inches")
            # print(f"Height: {height.inches:.2f} inches")

            # shape = add_rounded_rectangle(slide, left, top, width, height, text, color)
            shape = add_basic_shape(slide, shape_name, left, top, width, height, text, color, shape_size_ratio)
            shapes.append(shape)
    else:
        merge_start_time = time.time()
        for ocr_item in ocr_results:
            ocr_text = ocr_item['text']
            # print(f"OCR text: {ocr_text}")
            merged_item = {'Inches_box': ocr_item['Inches_box'], 'text': ocr_text, 'colour': ocr_item['colour']}

            for idx, gpt_item in enumerate(physical_results):

                if idx in matched_gpt_indices:  # 跳过已匹配的 GPT 项
                    continue

                gpt_text = gpt_item['text']

                if similarity(ocr_text, gpt_text) > 0.5:  # 判断相似度
                    # merged_item['colour'] = gpt_item['colour']
                    merged_item['shape'] = gpt_item['shape']
                    matched_gpt_indices.add(idx)  # 标记为已匹配
                    # print(f"Matched: '{ocr_text}' and '{gpt_text}'")
                    break  # 找到匹配项后跳出循环

            merged_results.append(merged_item)


        merge_end_time = time.time()
        merge_elapsed_time = merge_end_time - merge_start_time
        # print("结果融合时间: ", merge_elapsed_time)
        # print("merged_results: ", merged_results)

        for merged_item in merged_results:
            Inches_box = merged_item['Inches_box']
            text = merged_item['text']
            color = get_color_string(merged_item)

            if recognize_shape:
                shape_name = get_shape_string(merged_item)
            else:
                shape_name = "Rounded Rectangle"

            left = Inches(Inches_box[0][0])
            top = Inches(Inches_box[0][1])
            width = Inches((Inches_box[2][0] - Inches_box[0][0]))
            height = Inches((Inches_box[2][1] - Inches_box[0][1]))

            # print(f"Left: {left.inches:.2f} inches")
            # print(f"Top: {top.inches:.2f} inches")
            # print(f"Width: {width.inches:.2f} inches")
            # print(f"Height: {height.inches:.2f} inches")

            # shape = add_rounded_rectangle(slide, left, top, width, height, text, color)
            shape = add_basic_shape(slide, shape_name, left, top, width, height, text, color, shape_size_ratio)
            shapes.append(shape)


        logical_start_time = time.time()
        for shape in shapes:
            shape_text = shape.text  # 获取 shape 的文本内容
            for logic in logical_results:
                from_text = logic['from']
                to_text = logic['to']

                # 判断 shape 的文本与 logical_results 中的 `from` 是否相似
                if similarity(shape_text, from_text) > 0.7:
                    # 查找是否存在与 logical_results 中 `to` 相似的 shape
                    for target_shape in shapes:
                        target_shape_text = target_shape.text
                        if similarity(target_shape_text, to_text) > 0.7:
                            # 找到匹配的 shape，添加连接器
                            from_center_x, from_center_y = get_center(shape)
                            target_center_x, target_center_y = get_center(target_shape)

                            # 判断相对位置并设置连接点
                            if from_center_x < target_center_x:  # A在B的左边
                                start_x = from_center_x + shape.width / 2  # 从A右边的中心点
                                start_y = from_center_y  # A的中心Y坐标
                                end_x = target_center_x - target_shape.width / 2  # B左边的中心点
                                end_y = target_center_y  # B的中心Y坐标
                            elif from_center_x > target_center_x:  # A在B的右边
                                start_x = from_center_x - shape.width / 2  # 从A左边的中心点
                                start_y = from_center_y  # A的中心Y坐标
                                end_x = target_center_x + target_shape.width / 2  # B右边的中心点
                                end_y = target_center_y  # B的中心Y坐标
                            elif from_center_y < target_center_y:  # A在B的上方
                                start_x = from_center_x  # A的中心X坐标
                                start_y = from_center_y + shape.height / 2  # 从A下边的中心点
                                end_x = target_center_x  # B的中心X坐标
                                end_y = target_center_y - target_shape.height / 2  # B上边的中心点
                            else:  # A在B的下方
                                start_x = from_center_x  # A的中心X坐标
                                start_y = from_center_y - shape.height / 2  # 从A上边的中心点
                                end_x = target_center_x  # B的中心X坐标
                                end_y = target_center_y + target_shape.height / 2  # B下边的中心点

                            if add_lines_between_shapes:
                                if line_style == 'curve':
                                    # 添加连接器
                                    connector = slide.shapes.add_connector(
                                        MSO_CONNECTOR.CURVE,
                                        int(start_x), int(start_y),  # 连接线的起始位置
                                        int(end_x), int(end_y)  # 连接线的结束位置
                                    )
                                    connector.line.end_arrowhead = True
                                    # print(f"Added curve connector from ({start_x}, {start_y}) to ({end_x}, {end_y})")

                                elif line_style == 'elbow':
                                    # 添加连接器
                                    connector = slide.shapes.add_connector(
                                        MSO_CONNECTOR.ELBOW,
                                        int(start_x), int(start_y),  # 连接线的起始位置
                                        int(end_x), int(end_y)  # 连接线的结束位置
                                    )
                                    connector.line.end_arrowhead = True
                                    # print(f"Added elbow connector from ({start_x}, {start_y}) to ({end_x}, {end_y})")

                                elif line_style == 'straight':
                                    # 添加连接器
                                    connector = slide.shapes.add_connector(
                                        MSO_CONNECTOR.STRAIGHT,
                                        int(start_x), int(start_y),  # 连接线的起始位置
                                        int(end_x), int(end_y)  # 连接线的结束位置
                                    )
                                    connector.line.end_arrowhead = True
                                    # print(f"Added straight connector from ({start_x}, {start_y}) to ({end_x}, {end_y})")

                            # print(f"Connector added from '{shape_text}' to '{target_shape_text}'")

        logical_end_time = time.time()
        logical_elapsed_time = logical_end_time - logical_start_time
        # print("逻辑处理时间: ", logical_elapsed_time)

    # Save the PowerPoint file
    ppt_save_path = os.path.join('../logs/', user_tag, 'output', 'ChartPPT.pptx')
    prs.save(ppt_save_path)