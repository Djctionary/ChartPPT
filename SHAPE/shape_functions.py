from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def shape_process(gpt_results):
    """

    :param gpt_results:
    :return: physical_results, logical_results
    """
    physical_results = []
    logical_results = []

    # 根据不同的字段将结果分别存储到不同列表中
    for result in gpt_results:
        if 'text' in result:  # 判断是否为物理结果
            physical_results.append(result)
        elif 'from' in result and 'to' in result:  # 判断是否为逻辑结果
            logical_results.append(result)

    return physical_results, logical_results

## Decision
shape_mapping = {
    # Basic shapes
    "Rounded Rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "Rectangle": MSO_SHAPE.RECTANGLE,
    "Diamond": MSO_SHAPE.DIAMOND,
    "Ellipse": MSO_SHAPE.OVAL,
    "Parallelogram": MSO_SHAPE.PARALLELOGRAM,  
    "Flowchart Document": MSO_SHAPE.FLOWCHART_DOCUMENT,  
    "Trapezoid": MSO_SHAPE.TRAPEZOID, 
    "Flowchart Magnetic Disk": MSO_SHAPE.FLOWCHART_MAGNETIC_DISK,
}


def add_basic_shape(slide, shape_name, left, top, width, height, text, color, shape_size_ratio, font_size=14):
    # Map the shape name to the corresponding MSO_SHAPE
    shape_type = shape_mapping.get(shape_name, MSO_SHAPE.RECTANGLE)  # Default to Rectangle if not found

    width *= (shape_size_ratio / 100)
    height *= (shape_size_ratio / 100)
    # 如果 shape_name 是 "Diamond"，将宽度和高度乘以 2
    if shape_name == "Diamond" or shape_name == "Parallelogram" or shape_name == "Trapezoid":
        width *= 1.7
        height *= 1.7

    # Add the shape to the slide
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    
    # Set shape color
    shape.fill.solid()
    r, g, b = map(int, color.strip("RGB()").split(", "))
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    
    # Set text and font size
    text_frame = shape.text_frame

    text_frame.clear()

    lines = text.split('\n')

    if lines:
        p = text_frame.paragraphs[0]  # 获取第一个段落
        p.text = lines[0]  # 设置第一行文本
        p.font.size = Pt(font_size)  # 设置字体大小

        if r == g == b == 0:  # 如果颜色是黑色
            font_color = "RGB(255, 255, 255)"  # 设置为白色字体
        else:
            font_color = "RGB(0, 0, 0)"  # 设置为黑色字体
        r_font, g_font, b_font = map(int, font_color.strip("RGB()").split(", "))
        p.font.color.rgb = RGBColor(r_font, g_font, b_font)

    # 处理后续行
    for line in lines[1:]:  # 从第二行开始
        p = text_frame.add_paragraph()  # 添加新的段落
        p.text = line  # 设置每一行文本
        p.font.size = Pt(font_size)  # 设置字体大小

        # 设置字体颜色
        if r == g == b == 0:  # 如果颜色是黑色
            font_color = "RGB(255, 255, 255)"  # 设置为白色字体
        else:
            font_color = "RGB(0, 0, 0)"  # 设置为黑色字体
        r_font, g_font, b_font = map(int, font_color.strip("RGB()").split(", "))
        p.font.color.rgb = RGBColor(r_font, g_font, b_font)

    # p = text_frame.paragraphs[0]
    # p.text = text
    # p.font.size = Pt(font_size) # for one line text
    
    # Set vertical alignment to center and text alignment to center
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
    
    return shape